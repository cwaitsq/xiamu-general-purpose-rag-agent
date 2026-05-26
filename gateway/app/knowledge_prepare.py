from __future__ import annotations

import csv
import hashlib
import json
import re
import shutil
import zipfile
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from .config import Settings, load_settings
from .media_preprocess import MediaPreprocessError, preprocess_media_file
from .rag_clients import HttpJsonError, OpenAICompatibleChatClient
from .tenant_paths import DEFAULT_TENANT_ID, TenantWorkspace, ensure_workspace_bootstrap, get_tenant_workspace
from shared.tenant_profile import DEFAULT_TENANT_PROFILE, TenantProfile, load_tenant_profile

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover - optional dependency
    PdfReader = None

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover - optional dependency
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None


ROOT = Path(__file__).resolve().parents[2]
KNOWLEDGE_ROOT = ROOT / "知识库"
RAW_DIR = KNOWLEDGE_ROOT / "原始资料"
DOCS_DIR = KNOWLEDGE_ROOT / "已整理知识"
AUTO_DOCS_DIR = DOCS_DIR / "自动整理"
OUTPUT_DIR = KNOWLEDGE_ROOT / "整理结果"
REPORT_FILE = OUTPUT_DIR / "知识整理报告.json"
PREVIEW_FILE = OUTPUT_DIR / "知识整理预览.md"

SUPPORTED_SUFFIXES = {".txt", ".md", ".docx", ".csv", ".json", ".pptx", ".xlsx", ".pdf", ".srt", ".vtt"}
PUBLISH_STATUSES = {"active", "draft", "inactive"}
FAQ_QUESTION_RE = re.compile(r"^(?:[-*]\s*)?(?:q|问)[：:]\s*(.+?)\s*$", re.IGNORECASE)
FAQ_ANSWER_RE = re.compile(r"^(?:[-*]\s*)?(?:a|答)[：:]\s*(.*)\s*$", re.IGNORECASE)
MARKDOWN_HEADING_RE = re.compile(r"^#{1,3}\s+(.+?)\s*$")
NUMBERED_HEADING_RE = re.compile(r"^(?:\d+[.)]|[一二三四五六七八九十]+[、.])\s*(.+?)\s*$")
INTERNAL_KEYWORDS = [
    "内部",
    "仅内部",
    "内部使用",
    "特批",
    "赔付",
    "审批",
    "底价",
    "底线",
    "折扣权限",
    "黑名单",
]
FAQ_KEYWORDS = ["faq", "常见问题", "高频问题", "问答", "q&a", "qa"]
PRODUCT_KEYWORDS = [
    "产品",
    "规格",
    "材质",
    "包装",
    "logo",
    "定制",
    "打样",
    "样品",
    "moq",
    "起订量",
]
POLICY_KEYWORDS = [
    "规则",
    "流程",
    "要求",
    "政策",
    "付款",
    "报价",
    "交期",
    "物流",
    "清关",
    "售后",
    "索赔",
    "订单",
]
VIDEO_SUFFIXES = {".mp4", ".mov", ".avi", ".mkv", ".wmv", ".flv", ".webm"}
AUDIO_SUFFIXES = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"}
LEGACY_AUTO_DIR_NAMES = {"自动整理"}


@dataclass
class RawMaterial:
    index: int
    file_name: str
    relative_path: str
    suffix: str
    title: str
    text: str
    modified_at: str


@dataclass
class PreparedKnowledgeDoc:
    raw_file: str
    output_file: str
    kb_id: str
    title: str
    category: str
    status: str
    visibility: str
    source: str
    mode: str
    body: str
    warnings: list[str] = field(default_factory=list)

    def to_markdown(self) -> str:
        lines = [
            "---",
            f"kb_id: {self.kb_id}",
            f"title: {self.title}",
            f"category: {self.category}",
            f"status: {self.status}",
            f"version: {datetime.now().strftime('%Y-%m-%d')}",
            f"visibility: {self.visibility}",
            f"source: {self.source}",
            "---",
            "",
            self.body.strip(),
            "",
        ]
        return "\n".join(lines)

    def to_summary(self) -> dict[str, Any]:
        return {
            "raw_file": self.raw_file,
            "output_file": self.output_file,
            "kb_id": self.kb_id,
            "title": self.title,
            "category": self.category,
            "status": self.status,
            "visibility": self.visibility,
            "mode": self.mode,
            "warnings": list(self.warnings),
        }


@dataclass
class KnowledgePrepareResult:
    status: str
    message: str
    validate_only: bool
    publish_status: str
    use_llm: bool
    raw_files_total: int
    prepared_docs_total: int
    skipped_files_total: int
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    docs: list[dict[str, Any]] = field(default_factory=list)
    output_dir: str | None = None
    report_file: str | None = None
    preview_file: str | None = None
    started_at: str | None = None
    finished_at: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


def run_knowledge_prepare(
    *,
    validate_only: bool = False,
    publish_status: str = "draft",
    use_llm: bool = False,
    settings: Settings | None = None,
    tenant_id: str = DEFAULT_TENANT_ID,
    target_file_path: str | None = None,
) -> KnowledgePrepareResult:
    started_at = _iso_now()
    normalized_status = publish_status.strip().lower()
    errors: list[str] = []
    warnings: list[str] = []
    active_settings = settings or load_settings()
    workspace = get_tenant_workspace(tenant_id, collection_base=active_settings.qmd_collection)
    ensure_workspace_bootstrap(workspace)

    if normalized_status not in PUBLISH_STATUSES:
        errors.append(f"publish_status 不支持：{publish_status}")

    raw_materials, load_errors, load_warnings = load_raw_materials(
        settings=active_settings,
        raw_dir=workspace.raw_dir,
        target_file_path=target_file_path,
    )
    errors.extend(load_errors)
    warnings.extend(load_warnings)

    prepared_docs: list[PreparedKnowledgeDoc] = []
    if not errors:
        for material in raw_materials:
            doc = prepare_doc(
                material,
                publish_status=normalized_status,
                use_llm=use_llm,
                settings=active_settings,
            )
            prepared_docs.append(doc)
            warnings.extend(doc.warnings)

    result = KnowledgePrepareResult(
        status="failed" if errors else "success",
        message="原始资料校验通过" if not errors else "原始资料整理失败",
        validate_only=validate_only,
        publish_status=normalized_status,
        use_llm=use_llm,
        raw_files_total=len(raw_materials),
        prepared_docs_total=len(prepared_docs),
        skipped_files_total=max(0, len(raw_materials) - len(prepared_docs)),
        errors=errors,
        warnings=warnings,
        docs=[doc.to_summary() for doc in prepared_docs],
        output_dir=str(workspace.auto_docs_dir),
        report_file=str(workspace.prepare_report_file),
        preview_file=str(workspace.prepare_preview_file),
        started_at=started_at,
    )

    if not errors and not validate_only:
        write_prepared_docs(prepared_docs, workspace=workspace, replace_all=target_file_path is None)
        write_preview(prepared_docs, workspace=workspace)
        result.message = "原始资料已整理成标准知识文档"

    result.finished_at = _iso_now()
    write_report(result, workspace=workspace)
    return result


def load_raw_materials_with_settings(
    *,
    settings: Settings,
    raw_dir: Path,
    target_file_path: str | None = None,
) -> tuple[list[RawMaterial], list[str], list[str]]:
    errors: list[str] = []
    warnings: list[str] = []

    if not raw_dir.exists():
        return [], [f"未找到原始资料目录：{raw_dir}"], warnings

    if target_file_path:
        target_path = Path(target_file_path)
        if not target_path.exists():
            return [], [f"未找到指定文件：{target_file_path}"], warnings
        raw_paths = [target_path]
    else:
        raw_paths = sorted(path for path in raw_dir.rglob("*") if path.is_file())
    supported_paths = [path for path in raw_paths if path.suffix.lower() in SUPPORTED_SUFFIXES]
    unsupported_paths = [path for path in raw_paths if path.suffix.lower() not in SUPPORTED_SUFFIXES]

    if not supported_paths and not unsupported_paths:
        return [], [f"原始资料目录里没有可处理文件：{raw_dir}"], warnings

    materials: list[RawMaterial] = []
    material_index = 0
    for index, path in enumerate(supported_paths, start=1):
        try:
            text = read_raw_text(path)
        except (KeyError, OSError, ValueError, zipfile.BadZipFile) as exc:
            text = try_read_with_media_preprocess(path, settings=settings, warnings=warnings, errors=errors)
            if not text:
                errors.append(f"{path.name} 读取失败：{exc}")
                continue
        cleaned = normalize_text(text)
        if len(cleaned) < 20:
            errors.append(f"{path.name} 可提取正文过少，暂时无法整理")
            continue

        title = infer_title(path, cleaned)
        material_index += 1
        try:
            relative_path = path.relative_to(raw_dir).as_posix()
        except ValueError:
            relative_path = path.name
        materials.append(
            RawMaterial(
                index=material_index,
                file_name=path.name,
                relative_path=relative_path,
                suffix=path.suffix.lower(),
                title=title,
                text=cleaned,
                modified_at=datetime.fromtimestamp(path.stat().st_mtime).isoformat(timespec="seconds"),
            )
        )

    for path in unsupported_paths:
        preprocessed_text = try_preprocess_unsupported_path(path, settings=settings, warnings=warnings, errors=errors)
        if not preprocessed_text:
            warnings.append(build_unsupported_warning(path))
            continue
        cleaned = normalize_text(preprocessed_text)
        if len(cleaned) < 20:
            errors.append(f"{path.name} 预处理后正文过少，暂时无法整理")
            continue
        material_index += 1
        try:
            relative_path = path.relative_to(raw_dir).as_posix()
        except ValueError:
            relative_path = path.name
        materials.append(
            RawMaterial(
                index=material_index,
                file_name=path.name,
                relative_path=relative_path,
                suffix=path.suffix.lower(),
                title=normalize_title(path.stem),
                text=cleaned,
                modified_at=datetime.fromtimestamp(path.stat().st_mtime).isoformat(timespec="seconds"),
            )
        )

    if not materials:
        errors.append(f"原始资料目录里没有可处理文件：{raw_dir}")
    return materials, errors, warnings


def load_raw_materials(
    *,
    settings: Settings | None = None,
    raw_dir: Path | None = None,
    target_file_path: str | None = None,
) -> tuple[list[RawMaterial], list[str], list[str]]:
    active_settings = settings or load_settings()
    return load_raw_materials_with_settings(
        settings=active_settings,
        raw_dir=raw_dir or RAW_DIR,
        target_file_path=target_file_path,
    )


def try_read_with_media_preprocess(
    path: Path,
    *,
    settings: Settings,
    warnings: list[str],
    errors: list[str],
) -> str:
    if path.suffix.lower() != ".pdf":
        return ""
    try:
        result = preprocess_media_file(path, settings)
    except MediaPreprocessError as exc:
        warnings.append(f"{path.name} 需要 OCR，但当前预处理失败：{exc}")
        return ""
    warnings.extend(result.warnings)
    warnings.append(f"{path.name} 已按扫描版 PDF 做 OCR 预处理")
    return result.text


def try_preprocess_unsupported_path(
    path: Path,
    *,
    settings: Settings,
    warnings: list[str],
    errors: list[str],
) -> str:
    if path.suffix.lower() not in VIDEO_SUFFIXES | AUDIO_SUFFIXES | IMAGE_SUFFIXES:
        return ""
    try:
        result = preprocess_media_file(path, settings)
    except MediaPreprocessError as exc:
        warnings.append(f"{path.name} 预处理失败：{exc}")
        return ""
    warnings.extend(result.warnings)
    warnings.append(f"{path.name} 已完成媒体预处理，输出：{result.output_file}")
    return result.text


def read_raw_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return read_text_with_fallback(path)
    if suffix == ".docx":
        return read_docx_text(path)
    if suffix == ".csv":
        return read_csv_text(path)
    if suffix == ".json":
        return read_json_text(path)
    if suffix == ".pptx":
        return read_pptx_text(path)
    if suffix == ".xlsx":
        return read_xlsx_text(path)
    if suffix == ".pdf":
        return read_pdf_text(path)
    if suffix in {".srt", ".vtt"}:
        return read_subtitle_text(path)
    raise ValueError(f"不支持的文件类型：{path.suffix}")


def read_text_with_fallback(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def read_docx_text(path: Path) -> str:
    with zipfile.ZipFile(path) as archive:
        xml_bytes = archive.read("word/document.xml")
    root = ElementTree.fromstring(xml_bytes)
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: list[str] = []
    for paragraph in root.findall(".//w:p", namespace):
        fragments = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
        text = "".join(fragments).strip()
        if text:
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def read_csv_text(path: Path) -> str:
    rows: list[str] = []
    reader = csv.reader(read_text_with_fallback(path).splitlines())
    for row_index, row in enumerate(reader, start=1):
        normalized_row = [cell.strip() for cell in row if cell.strip()]
        if not normalized_row:
            continue
        rows.append(f"第 {row_index} 行：{' | '.join(normalized_row)}")
    return "\n".join(rows)


def read_json_text(path: Path) -> str:
    raw = read_text_with_fallback(path)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        return raw

    lines: list[str] = []
    flatten_json(data, lines=lines)
    return "\n".join(lines) if lines else raw


def read_pptx_text(path: Path) -> str:
    if Presentation is None:
        raise ValueError("当前环境没有安装 python-pptx，暂时不能处理 pptx")

    presentation = Presentation(path)
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        fragments: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            cleaned = normalize_text(text)
            if cleaned:
                fragments.append(cleaned)
        if fragments:
            slides.append(f"第 {slide_index} 页\n" + "\n".join(fragments))
    return "\n\n".join(slides)


def read_xlsx_text(path: Path) -> str:
    if load_workbook is None:
        raise ValueError("当前环境没有安装 openpyxl，暂时不能处理 xlsx")

    workbook = load_workbook(filename=path, data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        sheet_lines: list[str] = []
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if not values:
                continue
            sheet_lines.append(f"第 {row_index} 行：{' | '.join(values)}")
        if sheet_lines:
            lines.append(f"工作表：{sheet.title}\n" + "\n".join(sheet_lines))
    workbook.close()
    return "\n\n".join(lines)


def read_pdf_text(path: Path) -> str:
    if PdfReader is None:
        raise ValueError("当前环境没有安装 pypdf，暂时不能处理 pdf")

    reader = PdfReader(str(path))
    pages: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        if text:
            pages.append(f"第 {page_index} 页\n{text}")
    if pages:
        return "\n\n".join(pages)
    raise ValueError("PDF 没抽到有效正文，可能是扫描件，需要先做 OCR")


def read_subtitle_text(path: Path) -> str:
    text = read_text_with_fallback(path)
    lines: list[str] = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if re.fullmatch(r"\d+", line):
            continue
        if "-->" in line:
            continue
        if line.upper().startswith("WEBVTT"):
            continue
        lines.append(line)
    return "\n".join(lines)


def flatten_json(data: Any, *, prefix: str = "", lines: list[str]) -> None:
    if isinstance(data, dict):
        for key, value in data.items():
            next_prefix = f"{prefix}.{key}" if prefix else str(key)
            flatten_json(value, prefix=next_prefix, lines=lines)
        return
    if isinstance(data, list):
        for index, value in enumerate(data, start=1):
            next_prefix = f"{prefix}[{index}]"
            flatten_json(value, prefix=next_prefix, lines=lines)
        return
    value = "" if data is None else str(data).strip()
    if value:
        lines.append(f"{prefix}: {value}")


def build_unsupported_warning(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in VIDEO_SUFFIXES:
        return f"{path.name} 是视频，不能直接入库；请先做转写、字幕抽取、截图 OCR，再产出 txt/md/json。"
    if suffix in AUDIO_SUFFIXES:
        return f"{path.name} 是音频，不能直接入库；请先做语音转文字，再产出 txt/md/json。"
    if suffix in IMAGE_SUFFIXES:
        return f"{path.name} 是图片，不能直接入库；请先做 OCR，再产出 txt/md/json。"
    return f"暂不支持的文件类型，已跳过：{path.name}"


def prepare_doc(
    material: RawMaterial,
    *,
    publish_status: str,
    use_llm: bool,
    settings: Settings,
) -> PreparedKnowledgeDoc:
    if use_llm:
        llm_doc = try_prepare_with_llm(material, publish_status=publish_status, settings=settings)
        if llm_doc is not None:
            return llm_doc

    faq_pairs = extract_faq_pairs(material.text)
    category = detect_category(material, faq_pairs=faq_pairs)
    visibility = detect_visibility(material)
    title = material.title
    warnings: list[str] = []

    if faq_pairs:
        body = render_faq_body(title, faq_pairs)
    else:
        sections = extract_sections(material.text)
        if not sections:
            sections = [("正文", material.text)]
            warnings.append(f"{material.file_name} 没识别到明显结构，已按整篇正文整理")
        body = render_section_body(title, sections)

    file_name = build_output_file_name(material)
    kb_id = build_kb_id(material)
    return PreparedKnowledgeDoc(
        raw_file=material.relative_path,
        output_file=file_name,
        kb_id=kb_id,
        title=title,
        category=category,
        status=publish_status,
        visibility=visibility,
        source=f"raw::{material.relative_path}",
        mode="rule",
        body=body,
        warnings=warnings,
    )


def try_prepare_with_llm(
    material: RawMaterial,
    *,
    publish_status: str,
    settings: Settings,
) -> PreparedKnowledgeDoc | None:
    chat_client = OpenAICompatibleChatClient(settings)
    if not chat_client.is_configured():
        return None

    truncated_text = material.text[:12000]
    system_prompt = (
        "你是企业知识整理助手。"
        "你的任务是把原始资料整理成标准知识文档。"
        "只能基于原文整理，不允许补充原文没有的事实。"
        "信息不清楚就保留原话，不要猜。"
        "只输出 JSON，不要输出解释。"
    )
    user_prompt = (
        "请把下面原始资料整理成 JSON，字段要求如下：\n"
        "{\n"
        '  "title": "文档标题",\n'
        '  "category": "faq|policy|product",\n'
        '  "visibility": "external|internal",\n'
        '  "faqs": [{"question": "", "answer": ""}],\n'
        '  "sections": [{"title": "", "content": ""}]\n'
        "}\n"
        "如果更像问答，就填 faqs；如果更像说明文档，就填 sections。\n"
        "原始资料如下：\n"
        f"{truncated_text}"
    )
    try:
        content = chat_client.complete(system_prompt=system_prompt, user_prompt=user_prompt)
    except HttpJsonError:
        return None

    parsed = parse_json_object(content)
    if not isinstance(parsed, dict):
        return None

    title = normalize_title(str(parsed.get("title", "")).strip()) or material.title
    category = str(parsed.get("category", "")).strip().lower()
    if category not in {"faq", "policy", "product"}:
        category = detect_category(material, faq_pairs=extract_faq_pairs(material.text))

    visibility = str(parsed.get("visibility", "")).strip().lower()
    if visibility not in {"external", "internal"}:
        visibility = detect_visibility(material)

    warnings: list[str] = ["本篇使用 LLM 做了结构整理，建议人工抽查"]
    faqs = parsed.get("faqs", [])
    sections = parsed.get("sections", [])
    body = ""
    if isinstance(faqs, list):
        normalized_faqs: list[tuple[str, str]] = []
        for item in faqs:
            if not isinstance(item, dict):
                continue
            question = normalize_text(str(item.get("question", "")))
            answer = normalize_text(str(item.get("answer", "")))
            if question and answer:
                normalized_faqs.append((question, answer))
        if normalized_faqs:
            body = render_faq_body(title, normalized_faqs)
            category = "faq"

    if not body and isinstance(sections, list):
        normalized_sections: list[tuple[str, str]] = []
        for item in sections:
            if not isinstance(item, dict):
                continue
            section_title = normalize_title(str(item.get("title", "")))
            section_content = normalize_text(str(item.get("content", "")))
            if section_title and section_content:
                normalized_sections.append((section_title, section_content))
        if normalized_sections:
            body = render_section_body(title, normalized_sections)

    if not body:
        return None

    return PreparedKnowledgeDoc(
        raw_file=material.relative_path,
        output_file=build_output_file_name(material),
        kb_id=build_kb_id(material),
        title=title,
        category=category,
        status=publish_status,
        visibility=visibility,
        source=f"raw::{material.relative_path}",
        mode="llm",
        body=body,
        warnings=warnings,
    )


def parse_json_object(text: str) -> Any:
    content = text.strip()
    if not content:
        return None
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        pass

    start = content.find("{")
    end = content.rfind("}")
    if start < 0 or end <= start:
        return None
    try:
        return json.loads(content[start : end + 1])
    except json.JSONDecodeError:
        return None


def infer_title(path: Path, text: str) -> str:
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        heading_match = MARKDOWN_HEADING_RE.match(line)
        if heading_match:
            return normalize_title(heading_match.group(1))
        if len(line) <= 32 and not line.endswith(("。", "；", "：", ":", "，")):
            return normalize_title(line)
        break
    return normalize_title(path.stem)


def normalize_title(title: str) -> str:
    compact = re.sub(r"\s+", " ", title).strip()
    return compact[:60] if compact else ""


def normalize_text(text: str) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n").replace("\ufeff", "")
    normalized = re.sub(r"[ \t]+\n", "\n", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized.strip()


def detect_visibility(material: RawMaterial) -> str:
    haystack = f"{material.file_name}\n{material.text[:4000]}".lower()
    for keyword in INTERNAL_KEYWORDS:
        if keyword.lower() in haystack:
            return "internal"
    return "external"


def detect_category(material: RawMaterial, *, faq_pairs: list[tuple[str, str]]) -> str:
    haystack = f"{material.file_name}\n{material.title}\n{material.text[:4000]}".lower()
    if faq_pairs or any(keyword.lower() in haystack for keyword in FAQ_KEYWORDS):
        return "faq"

    product_hits = sum(1 for keyword in PRODUCT_KEYWORDS if keyword.lower() in haystack)
    policy_hits = sum(1 for keyword in POLICY_KEYWORDS if keyword.lower() in haystack)
    return "product" if product_hits > policy_hits else "policy"


def extract_faq_pairs(text: str) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    current_question = ""
    current_answer_lines: list[str] = []

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        question_match = FAQ_QUESTION_RE.match(line)
        if question_match:
            if current_question and current_answer_lines:
                answer = normalize_text("\n".join(current_answer_lines))
                if answer:
                    pairs.append((current_question, answer))
            current_question = question_match.group(1).strip()
            current_answer_lines = []
            continue

        answer_match = FAQ_ANSWER_RE.match(line)
        if answer_match and current_question:
            answer_text = answer_match.group(1).strip()
            if answer_text:
                current_answer_lines.append(answer_text)
            continue

        if current_question:
            current_answer_lines.append(line)

    if current_question and current_answer_lines:
        answer = normalize_text("\n".join(current_answer_lines))
        if answer:
            pairs.append((current_question, answer))
    return pairs


def extract_sections(text: str) -> list[tuple[str, str]]:
    markdown_sections = extract_markdown_sections(text)
    if markdown_sections:
        return markdown_sections

    blocks = [normalize_text(block) for block in re.split(r"\n\s*\n", text) if normalize_text(block)]
    sections: list[tuple[str, str]] = []
    for index, block in enumerate(blocks, start=1):
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue

        first_line = lines[0]
        numbered_match = NUMBERED_HEADING_RE.match(first_line)
        if numbered_match and len(lines) > 1:
            sections.append((normalize_title(numbered_match.group(1)), normalize_text("\n".join(lines[1:]))))
            continue

        if "：" in first_line:
            prefix, remainder = first_line.split("：", 1)
            if 1 <= len(prefix.strip()) <= 20 and remainder.strip():
                body_lines = [remainder.strip(), *lines[1:]]
                sections.append((normalize_title(prefix), normalize_text("\n".join(body_lines))))
                continue

        if ":" in first_line:
            prefix, remainder = first_line.split(":", 1)
            if 1 <= len(prefix.strip()) <= 20 and remainder.strip():
                body_lines = [remainder.strip(), *lines[1:]]
                sections.append((normalize_title(prefix), normalize_text("\n".join(body_lines))))
                continue

        if len(first_line) <= 24 and len(lines) > 1:
            sections.append((normalize_title(first_line), normalize_text("\n".join(lines[1:]))))
            continue

        sections.append((f"要点{index}", block))
    return [(title, body) for title, body in sections if title and body]


def extract_markdown_sections(text: str) -> list[tuple[str, str]]:
    sections: list[tuple[str, str]] = []
    current_title = ""
    current_lines: list[str] = []

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        heading_match = MARKDOWN_HEADING_RE.match(line.strip())
        if heading_match:
            heading = normalize_title(heading_match.group(1))
            if current_title and current_lines:
                sections.append((current_title, normalize_text("\n".join(current_lines))))
            current_title = heading
            current_lines = []
            continue
        if current_title:
            current_lines.append(line)

    if current_title and current_lines:
        sections.append((current_title, normalize_text("\n".join(current_lines))))
    return [(title, body) for title, body in sections if title and body]


def render_faq_body(title: str, faq_pairs: list[tuple[str, str]]) -> str:
    lines = [f"# {title}", ""]
    for question, answer in faq_pairs:
        lines.extend([f"## 问：{question}", "", f"答：{answer}", ""])
    return "\n".join(lines).strip()


def render_section_body(title: str, sections: list[tuple[str, str]]) -> str:
    lines = [f"# {title}", ""]
    for section_title, content in sections:
        lines.extend([f"## {section_title}", "", content, ""])
    return "\n".join(lines).strip()


def build_output_file_name(material: RawMaterial) -> str:
    digest = hashlib.sha1(material.relative_path.encode("utf-8")).hexdigest()[:10]
    return f"auto_{digest}.md"


def build_kb_id(material: RawMaterial) -> str:
    digest = hashlib.sha1(material.relative_path.encode("utf-8")).hexdigest()[:10]
    return f"auto_raw_{digest}"


def sanitize_file_name(name: str) -> str:
    sanitized = re.sub(r"[<>:\"/\\\\|?*]+", "_", name).strip(" .")
    return sanitized[:48] or "未命名知识"


def write_prepared_docs(
    docs: list[PreparedKnowledgeDoc],
    *,
    workspace: TenantWorkspace,
    replace_all: bool,
) -> None:
    if replace_all and workspace.auto_docs_dir.exists():
        safe_rmtree(workspace.auto_docs_dir, docs_dir=workspace.docs_dir)
    if replace_all:
        for legacy_dir_name in LEGACY_AUTO_DIR_NAMES:
            legacy_dir = workspace.docs_dir / legacy_dir_name
            if legacy_dir.exists():
                safe_rmtree(legacy_dir, docs_dir=workspace.docs_dir)
    workspace.auto_docs_dir.mkdir(parents=True, exist_ok=True)

    for doc in docs:
        output_path = workspace.auto_docs_dir / doc.output_file
        output_path.write_text(doc.to_markdown(), encoding="utf-8")


def safe_rmtree(path: Path, *, docs_dir: Path) -> None:
    resolved = path.resolve()
    allowed_root = docs_dir.resolve()
    if resolved == allowed_root or allowed_root not in resolved.parents:
        raise ValueError(f"不允许删除目录：{path}")
    shutil.rmtree(resolved)


def write_preview(docs: list[PreparedKnowledgeDoc], *, workspace: TenantWorkspace) -> None:
    workspace.prepare_output_dir.mkdir(parents=True, exist_ok=True)
    lines = ["# 知识整理预览", ""]
    for doc in docs:
        lines.extend(
            [
                f"## {doc.title}",
                f"- 原始文件：{doc.raw_file}",
                f"- 生成文件：{doc.output_file}",
                f"- 分类：{doc.category}",
                f"- 可见范围：{doc.visibility}",
                f"- 发布状态：{doc.status}",
                f"- 整理方式：{doc.mode}",
                "",
                doc.body,
                "",
            ]
        )
    workspace.prepare_preview_file.write_text("\n".join(lines), encoding="utf-8")


def write_report(result: KnowledgePrepareResult, *, workspace: TenantWorkspace) -> None:
    workspace.prepare_output_dir.mkdir(parents=True, exist_ok=True)
    workspace.prepare_report_file.write_text(json.dumps(result.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")


def _iso_now() -> str:
    return datetime.now().isoformat(timespec="seconds")


def run_knowledge_prepare(
    *,
    validate_only: bool = False,
    publish_status: str = "draft",
    use_llm: bool = False,
    settings: Settings | None = None,
    tenant_id: str = DEFAULT_TENANT_ID,
    target_file_path: str | None = None,
) -> KnowledgePrepareResult:
    started_at = _iso_now()
    normalized_status = publish_status.strip().lower()
    errors: list[str] = []
    warnings: list[str] = []
    active_settings = settings or load_settings()
    profile = load_tenant_profile(tenant_id)
    workspace = get_tenant_workspace(tenant_id, collection_base=active_settings.qmd_collection)
    ensure_workspace_bootstrap(workspace)

    if normalized_status not in PUBLISH_STATUSES:
        errors.append(f"publish_status 不支持：{publish_status}")

    raw_materials, load_errors, load_warnings = load_raw_materials(
        settings=active_settings,
        raw_dir=workspace.raw_dir,
        target_file_path=target_file_path,
    )
    errors.extend(load_errors)
    warnings.extend(load_warnings)

    prepared_docs: list[PreparedKnowledgeDoc] = []
    if not errors:
        for material in raw_materials:
            doc = prepare_doc(
                material,
                publish_status=normalized_status,
                use_llm=use_llm,
                settings=active_settings,
                profile=profile,
            )
            prepared_docs.append(doc)
            warnings.extend(doc.warnings)

    result = KnowledgePrepareResult(
        status="failed" if errors else "success",
        message="原始材料校验通过" if not errors else "原始材料整理失败",
        validate_only=validate_only,
        publish_status=normalized_status,
        use_llm=use_llm,
        raw_files_total=len(raw_materials),
        prepared_docs_total=len(prepared_docs),
        skipped_files_total=max(0, len(raw_materials) - len(prepared_docs)),
        errors=errors,
        warnings=warnings,
        docs=[doc.to_summary() for doc in prepared_docs],
        output_dir=str(workspace.auto_docs_dir),
        report_file=str(workspace.prepare_report_file),
        preview_file=str(workspace.prepare_preview_file),
        started_at=started_at,
    )

    if not errors and not validate_only:
        write_prepared_docs(prepared_docs, workspace=workspace, replace_all=target_file_path is None)
        write_preview(prepared_docs, workspace=workspace)
        result.message = "原始材料已整理成标准知识文档。"

    result.finished_at = _iso_now()
    write_report(result, workspace=workspace)
    return result


def prepare_doc(
    material: RawMaterial,
    *,
    publish_status: str,
    use_llm: bool,
    settings: Settings,
    profile: TenantProfile | None = None,
) -> PreparedKnowledgeDoc:
    active_profile = profile or DEFAULT_TENANT_PROFILE
    if use_llm:
        llm_doc = try_prepare_with_llm(material, publish_status=publish_status, settings=settings, profile=active_profile)
        if llm_doc is not None:
            return llm_doc

    faq_pairs = extract_faq_pairs(material.text)
    category = detect_category(material, faq_pairs=faq_pairs, profile=active_profile)
    visibility = detect_visibility(material)
    title = material.title
    warnings: list[str] = []

    if faq_pairs:
        body = render_faq_body(title, faq_pairs)
    else:
        sections = extract_sections(material.text)
        if not sections:
            sections = [("正文", material.text)]
            warnings.append(f"{material.file_name} 没识别到明显结构，已按整篇正文整理。")
        body = render_section_body(title, sections)

    file_name = build_output_file_name(material)
    kb_id = build_kb_id(material)
    return PreparedKnowledgeDoc(
        raw_file=material.relative_path,
        output_file=file_name,
        kb_id=kb_id,
        title=title,
        category=category,
        status=publish_status,
        visibility=visibility,
        source=f"raw::{material.relative_path}",
        mode="rule",
        body=body,
        warnings=warnings,
    )


def try_prepare_with_llm(
    material: RawMaterial,
    *,
    publish_status: str,
    settings: Settings,
    profile: TenantProfile | None = None,
) -> PreparedKnowledgeDoc | None:
    active_profile = profile or DEFAULT_TENANT_PROFILE
    chat_client = OpenAICompatibleChatClient(settings)
    if not chat_client.is_configured():
        return None

    truncated_text = material.text[:12000]
    system_prompt = active_profile.render(active_profile.llm_system_prompt)
    allowed_scopes = " | ".join(active_profile.allowed_kb_scopes)
    user_prompt = (
        "请把下面原始材料整理成 JSON，字段要求如下：\n"
        "{\n"
        '  "title": "文档标题",\n'
        f'  "category": "{allowed_scopes}",\n'
        '  "visibility": "external|internal",\n'
        '  "faqs": [{"question": "", "answer": ""}],\n'
        '  "sections": [{"title": "", "content": ""}]\n'
        "}\n"
        "如果更像问答，就填 faqs；如果更像说明文档，就填 sections。\n"
        "原始材料如下：\n"
        f"{truncated_text}"
    )
    try:
        content = chat_client.complete(system_prompt=system_prompt, user_prompt=user_prompt)
    except HttpJsonError:
        return None

    parsed = parse_json_object(content)
    if not isinstance(parsed, dict):
        return None

    title = normalize_title(str(parsed.get("title", "")).strip()) or material.title
    category = str(parsed.get("category", "")).strip().lower()
    if category not in active_profile.allowed_kb_scopes:
        category = detect_category(material, faq_pairs=extract_faq_pairs(material.text), profile=active_profile)

    visibility = str(parsed.get("visibility", "")).strip().lower()
    if visibility not in {"external", "internal"}:
        visibility = detect_visibility(material)

    warnings: list[str] = ["本篇使用 LLM 做了结构整理，建议人工抽检。"]
    faqs = parsed.get("faqs", [])
    sections = parsed.get("sections", [])
    body = ""
    if isinstance(faqs, list):
        normalized_faqs: list[tuple[str, str]] = []
        for item in faqs:
            if not isinstance(item, dict):
                continue
            question = normalize_text(str(item.get("question", "")))
            answer = normalize_text(str(item.get("answer", "")))
            if question and answer:
                normalized_faqs.append((question, answer))
        if normalized_faqs:
            body = render_faq_body(title, normalized_faqs)
            category = active_profile.faq_category

    if not body and isinstance(sections, list):
        normalized_sections: list[tuple[str, str]] = []
        for item in sections:
            if not isinstance(item, dict):
                continue
            section_title = normalize_title(str(item.get("title", "")))
            section_content = normalize_text(str(item.get("content", "")))
            if section_title and section_content:
                normalized_sections.append((section_title, section_content))
        if normalized_sections:
            body = render_section_body(title, normalized_sections)

    if not body:
        return None

    return PreparedKnowledgeDoc(
        raw_file=material.relative_path,
        output_file=build_output_file_name(material),
        kb_id=build_kb_id(material),
        title=title,
        category=category,
        status=publish_status,
        visibility=visibility,
        source=f"raw::{material.relative_path}",
        mode="llm",
        body=body,
        warnings=warnings,
    )


def detect_category(
    material: RawMaterial,
    *,
    faq_pairs: list[tuple[str, str]],
    profile: TenantProfile | None = None,
) -> str:
    active_profile = profile or DEFAULT_TENANT_PROFILE
    haystack = f"{material.file_name}\n{material.title}\n{material.text[:4000]}".lower()
    if faq_pairs or any(keyword.lower() in haystack for keyword in active_profile.faq_keywords):
        return active_profile.faq_category

    scores: list[tuple[str, int]] = []
    for category in active_profile.non_faq_categories:
        terms = active_profile.category_terms(category)
        score = sum(1 for keyword in terms if keyword.lower() in haystack)
        scores.append((category, score))

    best_category, best_score = active_profile.default_content_category, -1
    for category, score in scores:
        if score > best_score:
            best_category, best_score = category, score
    return best_category if best_score > 0 else active_profile.default_content_category
