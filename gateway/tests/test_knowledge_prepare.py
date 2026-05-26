from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from openpyxl import Workbook
from pptx import Presentation

from app import knowledge_ingest, knowledge_prepare
from app.config import Settings
from app.media_preprocess import MediaPreprocessResult
from app.tenant_paths import TenantWorkspace
from shared.tenant_profile import TenantProfile


def make_workspace(root: Path, tenant_id: str = "unit_test") -> TenantWorkspace:
    docs_dir = root / "docs"
    auto_docs_dir = docs_dir / "auto"
    prepare_output_dir = root / "prepare_output"
    ingest_output_dir = root / "ingest_output"
    return TenantWorkspace(
        tenant_id=tenant_id,
        tenant_slug=tenant_id,
        root=root,
        raw_dir=root / "raw",
        docs_dir=docs_dir,
        auto_docs_dir=auto_docs_dir,
        prepare_output_dir=prepare_output_dir,
        prepare_report_file=prepare_output_dir / "knowledge_prepare_report.json",
        prepare_preview_file=prepare_output_dir / "knowledge_prepare_preview.md",
        ingest_output_dir=ingest_output_dir,
        chunks_file=ingest_output_dir / "chunks.jsonl",
        ingest_report_file=ingest_output_dir / "knowledge_ingest_report.json",
        ingest_preview_file=ingest_output_dir / "chunk_preview.md",
        collection_name=f"foreign_trade_kb__{tenant_id}",
    )


class KnowledgePrepareTests(unittest.TestCase):
    @staticmethod
    def make_settings(*, media_preprocess_enabled: bool = False, media_output_dir: str = ".") -> Settings:
        return Settings(
            gateway_mode="qmd",
            gateway_api_key="",
            qmd_command="",
            qmd_node_bin="node",
            qmd_cli_path="",
            qmd_collection="foreign_trade_kb",
            qmd_search_mode="search",
            qmd_timeout_seconds=60,
            qdrant_url="",
            qdrant_api_key="",
            qdrant_collection="",
            qdrant_distance="Cosine",
            qdrant_timeout_seconds=30,
            retrieval_score_threshold=0.35,
            embedding_base_url="",
            embedding_api_key="",
            embedding_model="Qwen3-Embedding-0.6B",
            embedding_dimensions=None,
            llm_base_url="",
            llm_api_key="",
            llm_model="",
            llm_temperature=0.2,
            llm_enabled=False,
            app_log_dir=".",
            media_preprocess_enabled=media_preprocess_enabled,
            multimodal_base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
            multimodal_api_key="test-key",
            multimodal_timeout_seconds=180,
            image_ocr_model="qwen-vl-ocr",
            audio_asr_model="qwen3-asr-flash",
            video_understand_model="qwen3.5-omni-plus",
            media_max_file_bytes=20 * 1024 * 1024,
            media_output_dir=media_output_dir,
        )

    def test_prepare_raw_txt_files_into_standard_docs(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            workspace = make_workspace(Path(tmp_dir))
            workspace.raw_dir.mkdir(parents=True, exist_ok=True)
            (workspace.raw_dir / "客户问答.txt").write_text(
                "\n".join(
                    [
                        "客户常见问题",
                        "",
                        "Q: 可以先拿样品吗？",
                        "A: 可以，样品费和快递费按实际确认。",
                        "",
                        "Q: 支持定制 logo 吗？",
                        "A: 支持，但要先确认设计稿和起订量。",
                    ]
                ),
                encoding="utf-8",
            )

            with patch.object(knowledge_prepare, "get_tenant_workspace", return_value=workspace), patch.object(
                knowledge_prepare, "ensure_workspace_bootstrap", return_value=None
            ):
                result = knowledge_prepare.run_knowledge_prepare(
                    validate_only=False,
                    publish_status="active",
                    use_llm=False,
                    tenant_id=workspace.tenant_id,
                )

            self.assertEqual(result.status, "success")
            generated_files = list(workspace.auto_docs_dir.glob("*.md"))
            self.assertEqual(len(generated_files), 1)
            generated_text = generated_files[0].read_text(encoding="utf-8")
            self.assertIn("category: faq", generated_text)
            self.assertIn("visibility: external", generated_text)
            self.assertIn("## 问：可以先拿样品吗？", generated_text)

    def test_ingest_reads_markdown_from_subdirectories(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            workspace = make_workspace(Path(tmp_dir))
            workspace.auto_docs_dir.mkdir(parents=True, exist_ok=True)
            (workspace.auto_docs_dir / "auto_001_test.md").write_text(
                "\n".join(
                    [
                        "---",
                        "kb_id: auto_raw_001_test",
                        "title: 测试知识",
                        "category: policy",
                        "status: active",
                        "version: 2026-05-22",
                        "visibility: external",
                        "source: raw::test.txt",
                        "---",
                        "",
                        "# 测试知识",
                        "",
                        "## 付款方式",
                        "",
                        "支持 T/T。",
                        "",
                    ]
                ),
                encoding="utf-8",
            )

            docs, errors, _warnings = knowledge_ingest.load_and_validate_docs(docs_dir=workspace.docs_dir)

            self.assertEqual(len(errors), 0)
            self.assertEqual(len(docs), 1)
            self.assertEqual(docs[0].file_name, "auto_001_test.md")

    def test_ingest_skips_legacy_auto_dir_when_new_auto_dir_exists(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            workspace = make_workspace(Path(tmp_dir))
            workspace.auto_docs_dir.mkdir(parents=True, exist_ok=True)
            legacy_auto_dir = workspace.docs_dir / "自动整理"
            legacy_auto_dir.mkdir(parents=True, exist_ok=True)

            new_doc = workspace.auto_docs_dir / "auto_new.md"
            old_doc = legacy_auto_dir / "legacy_old.md"
            content = "\n".join(
                [
                    "---",
                    "kb_id: kb_demo",
                    "title: 测试知识",
                    "category: policy",
                    "status: active",
                    "version: 2026-05-22",
                    "visibility: external",
                    "source: raw::test.txt",
                    "---",
                    "",
                    "# 测试知识",
                    "",
                    "## 付款方式",
                    "",
                    "支持 T/T。",
                    "",
                ]
            )
            new_doc.write_text(content, encoding="utf-8")
            old_doc.write_text(content.replace("kb_demo", "kb_old"), encoding="utf-8")

            docs, errors, _warnings = knowledge_ingest.load_and_validate_docs(docs_dir=workspace.docs_dir)

            self.assertEqual(len(errors), 0)
            self.assertEqual([doc.file_name for doc in docs], ["auto_new.md"])

    def test_prepare_supports_pptx_and_xlsx(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            workspace = make_workspace(Path(tmp_dir))
            workspace.raw_dir.mkdir(parents=True, exist_ok=True)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "报价规则"
            slide.placeholders[1].text = "付款方式：T/T\n报价前确认规格、数量、贸易条款"
            pptx_path = workspace.raw_dir / "报价规则.pptx"
            presentation.save(pptx_path)

            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "物流时效"
            worksheet.append(["国家", "时效"])
            worksheet.append(["美国", "25天"])
            worksheet.append(["德国", "30天"])
            xlsx_path = workspace.raw_dir / "物流时效.xlsx"
            workbook.save(xlsx_path)
            workbook.close()

            with patch.object(knowledge_prepare, "get_tenant_workspace", return_value=workspace), patch.object(
                knowledge_prepare, "ensure_workspace_bootstrap", return_value=None
            ):
                result = knowledge_prepare.run_knowledge_prepare(
                    validate_only=False,
                    publish_status="active",
                    use_llm=False,
                    tenant_id=workspace.tenant_id,
                )

            self.assertEqual(result.status, "success")
            self.assertEqual(result.prepared_docs_total, 2)
            generated_files = sorted(path.name for path in workspace.auto_docs_dir.glob("*.md"))
            self.assertEqual(len(generated_files), 2)

    def test_prepare_warns_for_video_files(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            workspace = make_workspace(Path(tmp_dir))
            workspace.raw_dir.mkdir(parents=True, exist_ok=True)
            (workspace.raw_dir / "产品培训.mp4").write_bytes(b"fake-video")

            _materials, errors, warnings = knowledge_prepare.load_raw_materials(
                settings=self.make_settings(media_preprocess_enabled=False),
                raw_dir=workspace.raw_dir,
            )

            self.assertEqual(len(errors), 1)
            self.assertTrue(any("视频" in warning for warning in warnings))

    def test_prepare_accepts_video_after_media_preprocess(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            workspace = make_workspace(Path(tmp_dir))
            workspace.raw_dir.mkdir(parents=True, exist_ok=True)
            (workspace.raw_dir / "产品培训.mp4").write_bytes(b"fake-video")

            fake_settings = self.make_settings(
                media_preprocess_enabled=True,
                media_output_dir=str(Path(tmp_dir) / "media_output"),
            )

            fake_result = MediaPreprocessResult(
                source_file="产品培训.mp4",
                media_type="video",
                text="今天培训主要讲外贸报价、MOQ、交期确认。",
                output_file=str(Path(tmp_dir) / "media_output" / "自动转写" / "产品培训.video.txt"),
                warnings=[],
            )

            with patch.object(knowledge_prepare, "get_tenant_workspace", return_value=workspace), patch.object(
                knowledge_prepare, "ensure_workspace_bootstrap", return_value=None
            ), patch.object(knowledge_prepare, "preprocess_media_file", return_value=fake_result):
                result = knowledge_prepare.run_knowledge_prepare(
                    validate_only=False,
                    publish_status="active",
                    use_llm=False,
                    settings=fake_settings,
                    tenant_id=workspace.tenant_id,
                )

            self.assertEqual(result.status, "success")
            generated_files = list(workspace.auto_docs_dir.glob("*.md"))
            self.assertEqual(len(generated_files), 1)
            generated_text = generated_files[0].read_text(encoding="utf-8")
            self.assertIn("外贸报价", generated_text)


    def test_prepare_respects_tenant_profile_category_mapping(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            workspace = make_workspace(Path(tmp_dir), tenant_id="custom_tenant")
            workspace.raw_dir.mkdir(parents=True, exist_ok=True)
            (workspace.raw_dir / "help.txt").write_text(
                "\n".join(
                    [
                        "Q: 如何申请售后？",
                        "A: 请联系人工客服。",
                    ]
                ),
                encoding="utf-8",
            )

            profile = TenantProfile(
                faq_category="help",
                allowed_kb_scopes=("help", "policy"),
                default_kb_scope=("help", "policy"),
                category_keywords={"help": ("售后",), "policy": ("付款",)},
            )

            with patch.object(knowledge_prepare, "get_tenant_workspace", return_value=workspace), patch.object(
                knowledge_prepare, "ensure_workspace_bootstrap", return_value=None
            ), patch.object(knowledge_prepare, "load_tenant_profile", return_value=profile):
                result = knowledge_prepare.run_knowledge_prepare(
                    validate_only=False,
                    publish_status="active",
                    use_llm=False,
                    tenant_id=workspace.tenant_id,
                )

            self.assertEqual(result.status, "success")
            generated_files = list(workspace.auto_docs_dir.glob("*.md"))
            self.assertEqual(len(generated_files), 1)
            generated_text = generated_files[0].read_text(encoding="utf-8")
            self.assertIn("category: help", generated_text)


if __name__ == "__main__":
    unittest.main()
